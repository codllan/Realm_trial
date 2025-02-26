##run command##  streamlit run c:/Users/Greenlight.Remote1/Documents/Programming.Files/ProgrammingProjects/realm-trial/Realm_trial/app.py [ARGUMENTS] --server.fileWatcherType=none
import streamlit as st
import pymongo
from sentence_transformers import SentenceTransformer
from PIL import Image
import io
import base64
from datetime import datetime
import requests
import uuid
from pptx import Presentation  # For PowerPoint generation
from pptx.util import Inches

# MongoDB connection
@st.cache_resource
def init_connection():
    connection_string = st.secrets["mongo"]["connection_string"]
    client = pymongo.MongoClient(connection_string, tls=True)
    return client

client = init_connection()
db = client["vector_db"]
collection = db["data_entries"]

# Load embedding model
@st.cache_resource
def load_model():
    return SentenceTransformer('all-MiniLM-L6-v2', device='cpu')

model = load_model()

# xAI API setup
XAI_API_URL = "https://api.x.ai/v1/chat/completions"
XAI_API_KEY = st.secrets["xai"]["api_key"]

# Function to vectorize text
def vectorize_text(text):
    return model.encode(text).tolist()

# Function to encode image to base64
def image_to_base64(image):
    buffered = io.BytesIO()
    image.save(buffered, format="PNG")
    return base64.b64encode(buffered.getvalue()).decode()

# Function to create a thumbnail from base64
def create_thumbnail(base64_string, size=(100, 100)):
    image_data = base64.b64decode(base64_string)
    img = Image.open(io.BytesIO(image_data))
    img.thumbnail(size)
    buffered = io.BytesIO()
    img.save(buffered, format="PNG")
    return base64.b64encode(buffered.getvalue()).decode()

# Function to export document to PowerPoint
def export_to_pptx(doc_id, subject, main_text, slides):
    prs = Presentation()
    
    # Title Slide: Subject and Timestamp
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title and Content layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = subject
    subtitle.text = f"Document ID: {doc_id}\nUploaded: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # Main Text Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Main Content"
    content.text = main_text
    
    # Slides with Images and Notes
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(4)
        
        # Add image
        img_data = base64.b64decode(slide_data["screenshot"])
        img_stream = io.BytesIO(img_data)
        slide.shapes.add_picture(img_stream, left, top, width, height)
        
        # Add notes below image
        notes_top = top + height + Inches(0.5)
        textbox = slide.shapes.add_textbox(left, notes_top, width, Inches(1))
        text_frame = textbox.text_frame
        text_frame.text = slide_data["notes"] if slide_data["notes"] else "No notes"
    
    # Save to a BytesIO buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Function to query xAI
def ask_xai(question, context_chunks):
    context = "\n".join([chunk["text"] for chunk in context_chunks])
    prompt = f"Based on the following data:\n\n{context}\n\nAnswer this question: {question}. Carefully consider all provided information."
    headers = {"Authorization": f"Bearer {XAI_API_KEY}", "Content-Type": "application/json"}
    payload = {"model": "grok-beta", "messages": [{"role": "system", "content": "You are a helpful assistant."}, {"role": "user", "content": prompt}], "max_tokens": 300, "temperature": 0.7}
    try:
        response = requests.post(XAI_API_URL, headers=headers, json=payload)
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"]
    except requests.exceptions.RequestException as e:
        st.error(f"xAI API Error: {e}")
        return "Sorry, I couldn’t process your question."

# Streamlit UI
st.title("Vector Data Collector with xAI Question-Answering")

# Tabs
tab1, tab2, tab3 = st.tabs(["Enter Data", "Ask a Question", "Previous Uploads"])

# --- Data Entry Tab ---
with tab1:
    # st.header("Enter Your Data")
    
    if "subject" not in st.session_state:
        st.session_state.subject = ""
    if "main_text" not in st.session_state:
        st.session_state.main_text = ""
    if "slide_count" not in st.session_state:
        st.session_state.slide_count = 1
    if "notes" not in st.session_state:
        st.session_state.notes = {}
    if "slides_draft" not in st.session_state:
        st.session_state.slides_draft = {}

    subject = st.text_input("Subject", value=st.session_state.subject, key="subject_input")
    main_text = st.text_area("Main Paragraph", value=st.session_state.main_text, height=200, key="main_text_input")
    st.subheader("Add Screenshots and Notes")

    slides = []
    for i in range(st.session_state.slide_count):
        st.write(f"Slide {i+1}")
        screenshot_key = f"screen_{i}"
        notes_key = f"notes_{i}"

        screenshot = st.file_uploader(f"Screenshot {i+1}", type=["png", "jpg", "jpeg"], key=screenshot_key)
        if screenshot:
            image = Image.open(screenshot)
            st.session_state.slides_draft[screenshot_key] = {"screenshot": image_to_base64(image), "name": screenshot.name}
        
        if screenshot_key in st.session_state.slides_draft:
            st.write(f"Uploaded: {st.session_state.slides_draft[screenshot_key]['name']}")
        
        if notes_key not in st.session_state.notes:
            st.session_state.notes[notes_key] = ""
        notes = st.text_input(f"Notes for Slide {i+1}", value=st.session_state.notes[notes_key], key=notes_key)

        if screenshot_key in st.session_state.slides_draft:
            slides.append({"screenshot": st.session_state.slides_draft[screenshot_key]["screenshot"], "notes": notes})
            if i == st.session_state.slide_count - 1:
                st.session_state.slide_count += 1

    if st.button("Upload"):
        if subject.strip() and main_text.strip():
            doc_id = str(uuid.uuid4())
            chunks = [
                {"type": "subject", "text": subject, "slides": []},
                {"type": "main_text", "text": main_text, "slides": []}
            ]
            for slide in slides:
                slide_text = slide["notes"] if slide["notes"] else "Slide with no notes"
                chunks.append({"type": "slide", "text": slide_text, "slides": [slide]})

            for chunk in chunks:
                embedding = vectorize_text(chunk["text"])
                doc = {
                    "doc_id": doc_id,
                    "subject": subject,
                    "original_text": main_text,
                    "chunk_type": chunk["type"],
                    "chunk_text": chunk["text"],
                    "embedding": embedding,
                    "slides": chunk["slides"],
                    "timestamp": datetime.utcnow()
                }
                collection.insert_one(doc)
            st.success(f"Data uploaded successfully! Document ID: {doc_id}")
            
            st.session_state.subject = ""
            st.session_state.main_text = ""
            st.session_state.notes = {}
            st.session_state.slides_draft = {}
            st.session_state.slide_count = 1
        else:
            st.error("Please enter both a subject and some text.")

# --- Question-Answering Tab ---
with tab2:
    # st.header("Ask a Question")
    question = st.text_input("Enter your question about the data")
    if st.button("Ask"):
        if question.strip():
            query_embedding = vectorize_text(question)
            st.write("Debug: Query Embedding Sample:", query_embedding[:5])
            pipeline = [
                {"$vectorSearch": {"index": "vector_index", "path": "embedding", "queryVector": query_embedding, "limit": 10, "numCandidates": 1000}},
                {"$project": {"doc_id": 1, "subject": 1, "original_text": 1, "chunk_type": 1, "chunk_text": 1, "slides": 1, "score": {"$meta": "vectorSearchScore"}}},
                {"$limit": 6}
            ]
            try:
                results = list(collection.aggregate(pipeline))
                st.write("Debug: Raw Vector Search Results:", results)
                if results:
                    context_chunks = [{"text": r["chunk_text"], "score": r["score"]} for r in results]
                    st.write("Debug: Context Chunks:", context_chunks)
                    answer = ask_xai(question, context_chunks)
                    st.write("**Answer:**", answer)
                    st.write("**Referenced Chunks and Images:**")
                    for i, result in enumerate(results, 1):
                        doc_id = result.get("doc_id", "N/A")
                        st.write(f"{i}. [Doc ID: {doc_id}] [{result['chunk_type']}] Subject: {result['subject']} - {result['chunk_text']} (Score: {result['score']})")
                        if "slides" in result and result["slides"]:
                            for slide in result["slides"]:
                                image_data = base64.b64decode(slide["screenshot"])
                                st.image(image_data, caption=slide["notes"], width=300)
                else:
                    st.write("No relevant data found to answer your question.")
            except pymongo.errors.OperationFailure as e:
                st.error(f"MongoDB Error: {e}")
        else:
            st.error("Please enter a question.")

# --- Previous Uploads Tab ---
with tab3:
    # st.header("Previous Uploads")
    
    try:
        pipeline = [
            {"$group": {
                "_id": "$doc_id",
                "subject": {"$first": "$subject"},
                "original_text": {"$first": "$original_text"},
                "timestamp": {"$max": "$timestamp"},
                "slides": {"$push": "$slides"}
            }},
            {"$sort": {"timestamp": -1}}
        ]
        documents = list(collection.aggregate(pipeline))
        
        if documents:
            cols = st.columns(3)
            for idx, doc in enumerate(documents):
                col = cols[idx % 3]
                with col:
                    st.markdown(
                        f"""
                        <div style='border: 1px solid #ddd; border-radius: 8px; padding: 10px; margin: 10px 0; background-color: #f9f9f9; box-shadow: 2px 2px 5px rgba(0,0,0,0.1);'>
                            <h4>{doc['subject']}</h4>
                            <p style='font-size: 12px; color: #666;'>Doc ID: {doc['_id']}</p>
                            <p style='font-size: 12px; color: #666;'>Uploaded: {doc['timestamp'].strftime('%Y-%m-%d %H:%M:%S')}</p>
                        """, 
                        unsafe_allow_html=True
                    )
                    
                    all_slides = [slide for sublist in doc['slides'] for slide in sublist if slide]
                    if all_slides:
                        for slide in all_slides:
                            thumbnail = create_thumbnail(slide["screenshot"])
                            st.image(f"data:image/png;base64,{thumbnail}", width=100)
                            st.caption(slide["notes"] if slide["notes"] else "No notes")
                    else:
                        st.write("No attachments")
                    
                    # Export to PowerPoint button
                    pptx_buffer = export_to_pptx(doc["_id"], doc["subject"], doc["original_text"], all_slides)
                    st.download_button(
                        label="Export to PowerPoint",
                        data=pptx_buffer,
                        file_name=f"{doc['subject']}_{doc['_id']}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    
                    st.markdown("</div>", unsafe_allow_html=True)
        else:
            st.write("No previous uploads found.")
    except pymongo.errors.PyMongoError as e:
        st.error(f"Error fetching uploads: {e}")

# Instructions
# st.sidebar.write("Run with: `streamlit run app.py --server.fileWatcherType=none`")
# st.sidebar.write("Ensure MongoDB Atlas vector_index and xAI API key are set up.")